#!/usr/bin/env python3
"""
Création d'une présentation PowerPoint qui reflète les étapes du notebook
air_bnb_TCHIKAYA_Kevin.ipynb

- Chaque slide = une étape importante du notebook
- On explique : ce qui est fait + pourquoi c'est important
- Possibilité d'ajouter les captures d'écran correspondantes (plots, tableaux, etc.)

⚠️ À adapter :
- Les chemins dans IMAGE_PATHS pour pointer vers tes captures (.png, .jpg, ...)
"""

from pptx import Presentation
from pptx.util import Inches

# À ADAPTER : chemins vers tes captures d'écran
IMAGE_PATHS = {
    "aperçu_donnees": "captures/01_apercu_donnees.png",
    "selection_colonnes": "captures/02_selection_colonnes.png",
    "nettoyage_prix": "captures/03_nettoyage_prix.png",
    "gestion_na": "captures/04_gestion_na.png",
    "encodage_cat": "captures/05_encodage_categoriel.png",
    "feature_eng": "captures/06_feature_engineering.png",
    "corr_heatmap": "captures/07_correlation_heatmap.png",
    "boxplot_price_roomtype": "captures/08_boxplot_prix_roomtype.png",
    "outliers_iqr": "captures/09_outliers_iqr.png",
    "split_train_test": "captures/10_split_train_test.png",
    "reg_simple": "captures/11_regression_simple.png",
    "reg_multiple": "captures/12_regression_multiple_residus.png",
    "rf_base": "captures/13_random_forest_base.png",
    "rf_importance": "captures/14_importances_features_rf.png",
    "overfitting": "captures/15_overfitting_rf.png",
}

def add_bullet_slide(prs, title_text, bullets, image_key=None):
    """
    Ajoute une diapo Titre + Contenu avec liste de puces.
    bullets : liste de str ou (texte, niveau)
    image_key : clé de IMAGE_PATHS pour insérer une capture en bas de slide
    """
    slide_layout = prs.slide_layouts[1]  # Titre + contenu
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = title_text
    tf = content.text_frame
    tf.clear()

    first = True
    for b in bullets:
        level = 0
        text = b
        if isinstance(b, tuple):
            text, level = b

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = text
        p.level = level

    # Ajout éventuel d'une image en bas de slide
    if image_key and image_key in IMAGE_PATHS and IMAGE_PATHS[image_key]:
        try:
            img_path = IMAGE_PATHS[image_key]
            # Position : un peu en dessous du texte
            left = Inches(0.5)
            top = Inches(4.2)
            width = Inches(8)
            slide.shapes.add_picture(img_path, left, top, width=width)
        except Exception as e:
            print(f"[⚠] Impossible d'ajouter l'image '{image_key}' ({IMAGE_PATHS[image_key]}): {e}")

    return slide


def create_airbnb_presentation(output_name="presentation_airbnb_kevin_tchikaya.pptx"):
    prs = Presentation()

    # =========================
    # SLIDE 1 – Titre
    # =========================
    slide_layout = prs.slide_layouts[0]  # Titre principal
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "IA - AirBnB"
    subtitle.text = (
        "Objectif : préparer les données AirBnB pour les exploiter avec une IA\n"
        "Prédiction du prix des logements\n\n"
        "Notebook : air_bnb_TCHIKAYA_Kevin.ipynb\n"
        "Auteur : Kevin TCHIKAYA"
    )

    # =========================
    # SLIDE 2 – Contexte & objectif du projet
    # =========================
    bullets_contexte = [
        "Contexte : prédire le prix d'un nouveau logement sur AirBnB.",
        "Question centrale : quel est le bon prix à proposer ?",
        "Enjeux pour le propriétaire :",
        ("Prix trop élevé → peu de réservations", 1),
        ("Prix trop bas → perte de revenus", 1),
        "Rôle de l'IA : proposer un prix cohérent basé sur l'historique des annonces.",
    ]
    add_bullet_slide(prs, "Contexte du projet & objectif", bullets_contexte)

    # =========================
    # SLIDE 3 – Chargement & aperçu des données
    # =========================
    bullets_data = [
        "Données AirBnB chargées depuis listings.csv (InsideAirBnB).",
        "Création d'un DataFrame pandas (df) et d'une version filtrée (df_selected).",
        "Vérification rapide : dimensions, premières lignes, types des colonnes.",
        "Pourquoi cette étape ?",
        ("Comprendre la structure brute des données avant de transformer quoi que ce soit.", 1),
        ("Repérer les colonnes vraiment utiles pour la suite.", 1),
    ]
    add_bullet_slide(prs, "Chargement & aperçu des données", bullets_data,
                     image_key="aperçu_donnees")

    # =========================
    # SLIDE 4 – Sélection des colonnes pertinentes
    # =========================
    bullets_selection = [
        "Sélection d'un sous-ensemble de colonnes :",
        ("price, room_type, neighbourhood_cleansed,", 1),
        ("minimum_nights, number_of_reviews,", 1),
        ("review_scores_rating, beds, bedrooms, accommodates", 1),
        "Objectif : se concentrer sur les variables qui influencent le prix.",
        "Pourquoi cette étape ?",
        ("Réduire le bruit et la complexité du dataset.", 1),
        ("Faciliter l'interprétation des modèles.", 1),
    ]
    add_bullet_slide(prs, "Prétraitement – Sélection des colonnes", bullets_selection,
                     image_key="selection_colonnes")

    # =========================
    # SLIDE 5 – Nettoyage de la variable prix
    # =========================
    bullets_prix = [
        "La colonne price est au format texte avec symbole '$' et virgules.",
        "Nettoyage effectué dans le notebook :",
        ("Suppression des caractères non numériques ( $, , ).", 1),
        ("Conversion en float pour pouvoir faire des calculs.", 1),
        "Pourquoi cette étape ?",
        ("Permet d'utiliser le prix comme vraie variable numérique.", 1),
        ("Indispensable pour la régression et les statistiques.", 1),
    ]
    add_bullet_slide(prs, "Nettoyage de la variable prix", bullets_prix,
                     image_key="nettoyage_prix")

    # =========================
    # SLIDE 6 – Conversion numérique & gestion des NaN
    # =========================
    bullets_na = [
        "Colonnes converties en numérique (to_numeric avec errors='coerce') :",
        ("review_scores_rating, minimum_nights, number_of_reviews,", 1),
        ("beds, bedrooms, accommodates", 1),
        "Gestion des valeurs manquantes :",
        ("review_scores_rating → remplacée par la moyenne.", 1),
        ("minimum_nights, number_of_reviews, beds, bedrooms, accommodates → médiane.", 1),
        "Pourquoi cette étape ?",
        ("Les modèles de ML n'acceptent pas les NaN.", 1),
        ("Limiter la perte d'information en évitant de supprimer trop de lignes.", 1),
    ]
    add_bullet_slide(prs, "Conversion numérique & valeurs manquantes", bullets_na,
                     image_key="gestion_na")

    # =========================
    # SLIDE 7 – Suppression des doublons & lignes invalides
    # =========================
    bullets_doublons = [
        "Suppression des doublons sur df_selected.",
        "Suppression des lignes sans valeur de price ou room_type.",
        "Objectif : garder uniquement des observations exploitables et uniques.",
        "Pourquoi cette étape ?",
        ("Éviter que certaines annonces comptent plusieurs fois.", 1),
        ("Améliorer la qualité de l'apprentissage et des statistiques.", 1),
    ]
    add_bullet_slide(prs, "Nettoyage – Doublons & lignes manquantes", bullets_doublons)

    # =========================
    # SLIDE 8 – Encodage des variables catégorielles
    # =========================
    bullets_cat = [
        "Variables catégorielles :",
        ("room_type (Entire home/apt, Private room, Shared room...)", 1),
        ("neighbourhood_cleansed (quartiers)", 1),
        "Encodage réalisé dans le notebook (LabelEncoder) :",
        ("room_type_encoded", 1),
        ("neighbourhood_encoded", 1),
        "Pourquoi cette étape ?",
        ("Transformer les catégories en nombres pour les modèles.", 1),
        ("Indispensable pour la régression linéaire et RandomForest.", 1),
    ]
    add_bullet_slide(prs, "Encodage des variables catégorielles", bullets_cat,
                     image_key="encodage_cat")

    # =========================
    # SLIDE 9 – Feature engineering
    # =========================
    bullets_fe = [
        "Création de nouvelles features à partir des colonnes existantes :",
        ("bed_bedroom_ratio = beds / bedrooms (si possible).", 1),
        ("size_category (petit / moyen / grand logement) en fonction de accommodates.", 1),
        ("size_category_encoded pour utiliser la catégorie dans les modèles.", 1),
        "Parfois aussi une variable de type popularity_score (ex : number_of_reviews).",
        "Pourquoi cette étape ?",
        ("Enrichir l'information donnée au modèle.", 1),
        ("Permet de mieux capturer la structure réelle des logements.", 1),
    ]
    add_bullet_slide(prs, "Feature engineering", bullets_fe,
                     image_key="feature_eng")

    # =========================
    # SLIDE 10 – Analyse descriptive & corrélations
    # =========================
    bullets_desc = [
        "Analyse de la distribution du prix et des variables numériques.",
        "Matrice de corrélation entre prix et les autres variables.",
        "Visualisations typiques dans le notebook :",
        ("Heatmap de corrélation.", 1),
        ("Boxplots du prix par type de chambre.", 1),
        "Pourquoi cette étape ?",
        ("Identifier quelles variables sont les plus liées au prix.", 1),
        ("Détecter des relations non linéaires ou des anomalies.", 1),
    ]
    add_bullet_slide(prs, "Analyse descriptive & corrélations", bullets_desc,
                     image_key="corr_heatmap")

    # =========================
    # SLIDE 11 – Gestion des valeurs aberrantes (outliers)
    # =========================
    bullets_outliers = [
        "Utilisation de la méthode IQR (InterQuartile Range) sur :",
        ("price, accommodates, number_of_reviews", 1),
        "Suppression des observations avec valeurs extrêmes.",
        "Pourquoi cette étape ?",
        ("Les outliers peuvent fausser la régression (coefficients, R²).", 1),
        ("Améliore la stabilité et la robustesse des modèles.", 1),
    ]
    add_bullet_slide(prs, "Gestion des valeurs aberrantes (IQR)", bullets_outliers,
                     image_key="outliers_iqr")

    # =========================
    # SLIDE 12 – Construction de X, y & split train/test
    # =========================
    bullets_split = [
        "Définition de la variable cible y = price.",
        "Définition de X avec les features :",
        ("accommodates, bedrooms, beds, review_scores_rating,", 1),
        ("number_of_reviews, minimum_nights,", 1),
        ("room_type_encoded, neighbourhood_encoded,", 1),
        ("bed_bedroom_ratio, size_category_encoded", 1),
        "Séparation train/test (ex : 60% / 40%).",
        "Pourquoi cette étape ?",
        ("Évaluer les modèles sur des données non vues.", 1),
        ("Mesurer le généralisation, pas seulement l'apprentissage.", 1),
    ]
    add_bullet_slide(prs, "Construction de X, y & séparation train/test", bullets_split,
                     image_key="split_train_test")

    # =========================
    # SLIDE 13 – Régression linéaire simple
    # =========================
    bullets_reg_simple = [
        "Choix d'une seule variable explicative fortement corrélée au prix.",
        ("Exemple : accommodates (nombre de personnes acceptées).", 1),
        "Entraînement d'une régression linéaire simple : price ~ accommodates.",
        "Évaluation dans le notebook :",
        ("R², MAE, RMSE sur le jeu de test.", 1),
        "Pourquoi cette étape ?",
        ("Servir de baseline très simple.", 1),
        ("Comprendre l'effet isolé d'une variable sur le prix.", 1),
    ]
    add_bullet_slide(prs, "Modélisation – Régression linéaire simple", bullets_reg_simple,
                     image_key="reg_simple")

    # =========================
    # SLIDE 14 – Régression linéaire multiple
    # =========================
    bullets_reg_multiple = [
        "Utilisation de plusieurs variables explicatives :",
        ("Caractéristiques du logement + encodage des catégories.", 1),
        "Entraînement de la régression linéaire multiple sur X_train / y_train.",
        "Évaluation :",
        ("R², MAE, RMSE sur train et test.", 1),
        ("Analyse des résidus (erreurs) pour vérifier les hypothèses.", 1),
        "Pourquoi cette étape ?",
        ("Combiner plusieurs facteurs pour de meilleures prédictions.", 1),
        ("Comparer les performances à la régression simple.", 1),
    ]
    add_bullet_slide(prs, "Modélisation – Régression linéaire multiple", bullets_reg_multiple,
                     image_key="reg_multiple")

    # =========================
    # SLIDE 15 – Random Forest (modèle non linéaire)
    # =========================
    bullets_rf = [
        "Mise en place d'un modèle RandomForestRegressor.",
        "Objectif : capturer des relations non linéaires et des interactions.",
        "Comparaison avec la régression linéaire multiple :",
        ("Regarder R², MAE, RMSE.", 1),
        ("Analyse des importances de variables.", 1),
        "Pourquoi cette étape ?",
        ("Tester un modèle plus puissant sur les mêmes données.", 1),
        ("Voir quelles variables sont les plus importantes pour la prédiction.", 1),
    ]
    add_bullet_slide(prs, "Modélisation – Random Forest", bullets_rf,
                     image_key="rf_importance")

    # =========================
    # SLIDE 16 – Surapprentissage & conclusion modèles
    # =========================
    bullets_overfit = [
        "Vérification du surapprentissage (overfitting) :",
        ("Comparaison des scores entre train et test.", 1),
        ("Utilisation de la validation croisée (CV) pour la stabilité.", 1),
        "Analyse de la qualité du modèle :",
        ("Équilibre entre performance et généralisation.", 1),
        "Conclusion :",
        ("La régression multiple fournit une base interprétable.", 1),
        ("La Random Forest peut améliorer la performance au prix de la complexité.", 1),
    ]
    add_bullet_slide(prs, "Analyse du surapprentissage & conclusion", bullets_overfit,
                     image_key="overfitting")

    # =========================
    # Sauvegarde
    # =========================
    prs.save(output_name)
    print(f"✅ Présentation créée avec succès : {output_name}")


if __name__ == "__main__":
    try:
        create_airbnb_presentation()
    except ImportError:
        print("❌ Installe d'abord python-pptx :")
        print("   pip install python-pptx")
